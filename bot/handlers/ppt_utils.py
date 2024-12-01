from pptx import Presentation

def create_presentation(topic, content):
    presentation = Presentation()

    # Add title slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title = slide.shapes.title
    title.text = f"Presentation on {topic}"

    # Add content slides
    for i, slide_content in enumerate(content, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = f"Slide {i}"
        slide.placeholders[1].text = slide_content

    # Save presentation to file
    filename = f"{topic.replace(' ', '_')}.pptx"
    presentation.save(filename)
    return filename
